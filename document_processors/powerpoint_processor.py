"""
PowerPoint文档处理器
处理PPTX和PPT格式文档的文本提取和元数据提取
"""

import logging
from typing import List, Dict, Any
import os

from .base_processor import BaseDocumentProcessor

# 日志配置
logger = logging.getLogger(__name__)

class PowerPointProcessor(BaseDocumentProcessor):
    """PowerPoint文档处理器"""
    
    def __init__(self):
        """初始化PowerPoint处理器"""
        super().__init__()
    
    def can_process(self, file_path: str) -> bool:
        """判断是否可以处理PowerPoint文件"""
        return file_path.lower().endswith(('.pptx', '.ppt'))
    
    def get_supported_extensions(self) -> List[str]:
        """获取支持的文件扩展名"""
        return ['.pptx', '.ppt']
    
    def extract_text(self, file_path: str) -> str:
        """
        从PowerPoint文件中提取文本
        
        Args:
            file_path: PowerPoint文件路径
        
        Returns:
            提取的文本内容
        """
        try:
            # 目前主要支持.pptx格式
            if file_path.lower().endswith('.pptx'):
                return self._extract_from_pptx(file_path)
            else:
                # 对于.ppt文件，暂时提供提示信息
                logger.warning(f"暂不支持旧版PPT格式，请转换为PPTX格式: {file_path}")
                return f"提示：检测到旧版PPT文件 {os.path.basename(file_path)}，建议转换为PPTX格式以获得更好的文本提取效果。"
                
        except Exception as e:
            logger.error(f"从PowerPoint文档提取文本失败: {file_path} - {e}")
            return ""
    
    def _extract_from_pptx(self, file_path: str) -> str:
        """从PPTX文件提取文本"""
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            
            content = []
            content.append(f"PowerPoint演示文稿：{os.path.basename(file_path)}")
            content.append(f"幻灯片总数：{len(prs.slides)}")
            content.append("=" * 60)
            
            for i, slide in enumerate(prs.slides, 1):
                content.append(f"\n=== 幻灯片 {i} ===")
                
                slide_content = []
                
                # 提取幻灯片中的所有文本
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text = shape.text.strip()
                        if text:
                            # 判断是否是标题（通常是第一个文本框或较大的文本）
                            if shape.shape_type == 1:  # 文本框类型
                                # 简单的标题识别：通常标题较短且位置靠上
                                if len(text.split('\n')) == 1 and len(text) < 100:
                                    slide_content.insert(0, f"【标题】{text}")
                                else:
                                    slide_content.append(text)
                            else:
                                slide_content.append(text)
                    
                    # 提取表格内容
                    if hasattr(shape, "table"):
                        try:
                            table = shape.table
                            content.append("\n--- 幻灯片表格 ---")
                            for row in table.rows:
                                row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                                if row_text:
                                    slide_content.append(row_text)
                        except:
                            pass
                
                # 添加幻灯片内容
                if slide_content:
                    content.extend(slide_content)
                else:
                    content.append("（此幻灯片无文本内容）")
                
                content.append("")  # 空行分隔
            
            return "\n".join(content)
            
        except Exception as e:
            logger.error(f"从PPTX文件提取文本失败: {file_path} - {e}")
            return ""
    
    def extract_metadata(self, file_path: str) -> Dict[str, Any]:
        """
        提取PowerPoint文档元数据
        
        Args:
            file_path: PowerPoint文件路径
            
        Returns:
            PowerPoint元数据字典
        """
        try:
            metadata_dict = {
                'document_type': 'PowerPoint',
                'file_extension': os.path.splitext(file_path)[1].lower()
            }
            
            # 尝试从PPTX文件提取详细元数据
            if file_path.lower().endswith('.pptx'):
                try:
                    from pptx import Presentation
                    prs = Presentation(file_path)
                    
                    # 获取演示文稿基本信息
                    metadata_dict.update({
                        'slide_count': len(prs.slides),
                        'slide_width': prs.slide_width,
                        'slide_height': prs.slide_height
                    })
                    
                    # 获取核心属性
                    try:
                        core_props = prs.core_properties
                        metadata_dict.update({
                            'title': core_props.title or '',
                            'author': core_props.author or '',
                            'subject': core_props.subject or '',
                            'keywords': core_props.keywords or '',
                            'comments': core_props.comments or '',
                            'last_modified_by': core_props.last_modified_by or '',
                            'created': str(core_props.created) if core_props.created else '',
                            'modified': str(core_props.modified) if core_props.modified else '',
                            'category': core_props.category or ''
                        })
                    except Exception as e:
                        logger.debug(f"提取PPTX核心属性失败: {e}")
                    
                    # 分析幻灯片内容统计
                    total_shapes = 0
                    text_shapes = 0
                    total_text_length = 0
                    slides_with_text = 0
                    slides_with_images = 0
                    slides_with_tables = 0
                    
                    for slide in prs.slides:
                        slide_has_text = False
                        slide_has_image = False
                        slide_has_table = False
                        
                        for shape in slide.shapes:
                            total_shapes += 1
                            
                            # 统计文本
                            if hasattr(shape, "text") and shape.text.strip():
                                text_shapes += 1
                                total_text_length += len(shape.text)
                                slide_has_text = True
                            
                            # 统计图片
                            if shape.shape_type == 13:  # 图片类型
                                slide_has_image = True
                            
                            # 统计表格
                            if hasattr(shape, "table"):
                                slide_has_table = True
                        
                        if slide_has_text:
                            slides_with_text += 1
                        if slide_has_image:
                            slides_with_images += 1
                        if slide_has_table:
                            slides_with_tables += 1
                    
                    metadata_dict.update({
                        'total_shapes': total_shapes,
                        'text_shapes': text_shapes,
                        'total_text_length': total_text_length,
                        'slides_with_text': slides_with_text,
                        'slides_with_images': slides_with_images,
                        'slides_with_tables': slides_with_tables,
                        'avg_text_per_slide': total_text_length // max(len(prs.slides), 1),
                        'text_density': text_shapes / max(total_shapes, 1)
                    })
                    
                except Exception as e:
                    logger.warning(f"提取PPTX详细元数据失败: {file_path} - {e}")
            
            return metadata_dict
            
        except Exception as e:
            logger.error(f"提取PowerPoint元数据失败: {file_path} - {e}")
            return {
                'document_type': 'PowerPoint',
                'error': str(e)
            } 